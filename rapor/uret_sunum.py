"""
NetGuard — Kurumsal Sunum (Açık Tema, 10 slayt)
Tasarım ilkesi: Her slayt 5 saniyede anlaşılır — tek mesaj, büyük görsel, max 3 madde.
python3 uret_sunum.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DEST = "/home/mehmet/netguard/rapor/NetGuard_Sunum.pptx"

# ── Renk Paleti ───────────────────────────────────────────────────────────────
C_BG      = RGBColor(0xFF, 0xFF, 0xFF)
C_PANEL   = RGBColor(0xF1, 0xF5, 0xF9)  # slate-100
C_PANEL2  = RGBColor(0xE2, 0xE8, 0xF0)  # slate-200
C_NAVY    = RGBColor(0x0F, 0x17, 0x2A)  # başlık
C_DARK    = RGBColor(0x1E, 0x29, 0x3B)  # metin
C_MUTED   = RGBColor(0x64, 0x74, 0x8B)  # ikincil
C_BORDER  = RGBColor(0xCB, 0xD5, 0xE1)
C_BLUE    = RGBColor(0x02, 0x84, 0xC7)
C_BLUE_LT = RGBColor(0xE0, 0xF2, 0xFE)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_GREEN_LT= RGBColor(0xD1, 0xFA, 0xE5)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_AMBER_LT= RGBColor(0xFE, 0xF3, 0xC7)
C_VIOLET  = RGBColor(0x7C, 0x3A, 0xED)
C_VIO_LT  = RGBColor(0xED, 0xE9, 0xFE)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_RED_LT  = RGBColor(0xFE, 0xE2, 0xE2)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H  = Inches(7.5)

# ── Temel Yardımcılar ─────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = SLIDE_W
    p.slide_height = SLIDE_H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = color

def box(slide, x, y, w, h, fill=None, line=None, lw=1.0):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else:    s.line.fill.background()
    return s

def t(slide, text, x, y, w, h, size=14, bold=False, italic=False,
      color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color or C_DARK
    return tb

def line(slide, x, y, w, color=C_BORDER, thick=1.0):
    box(slide, x, y, w, Pt(thick), fill=color)

def header(slide, title, subtitle=None, accent=C_BLUE):
    """Üst şerit + başlık."""
    box(slide, 0, 0, SLIDE_W, Inches(0.07), fill=accent)
    box(slide, Inches(0.5), Inches(0.14), Pt(5), Inches(0.68), fill=accent)
    t(slide, title, Inches(0.65), Inches(0.1), Inches(11.5), Inches(0.65),
      size=30, bold=True, color=C_NAVY)
    if subtitle:
        t(slide, subtitle, Inches(0.65), Inches(0.75), Inches(11.5), Inches(0.3),
          size=13, color=C_MUTED)
    line(slide, Inches(0.5), Inches(1.1), Inches(12.3))

def photo_box(slide, x, y, w, h, label="Ekran görüntüsü eklenecek"):
    box(slide, x, y, w, h, fill=C_PANEL, line=C_BORDER, lw=1.0)
    t(slide, "📷", x, y + h/2 - Inches(0.55), w, Inches(0.5),
      size=28, color=C_BORDER, align=PP_ALIGN.CENTER)
    t(slide, f"[ {label} ]", x, y + h/2 - Inches(0.05), w, Inches(0.4),
      size=12, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

def kpi(slide, x, y, w, h, value, label, accent=C_BLUE, bg_c=None):
    box(slide, x, y, w, h, fill=bg_c or C_BLUE_LT, line=accent, lw=1.5)
    t(slide, value, x, y + Inches(0.06), w, h * 0.6,
      size=34, bold=True, color=accent, align=PP_ALIGN.CENTER)
    t(slide, label, x, y + h * 0.6, w, h * 0.35,
      size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

def bullet_row(slide, x, y, w, text, accent=C_BLUE, size=13):
    box(slide, x, y + Inches(0.09), Pt(9), Pt(9), fill=accent)
    t(slide, text, x + Inches(0.2), y, w - Inches(0.2), Inches(0.35),
      size=size, color=C_DARK)

def bottom_bar(slide, text, color=C_NAVY):
    box(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), fill=color)
    t(slide, text, Inches(0.4), SLIDE_H - Inches(0.35), SLIDE_W - Inches(0.8), Inches(0.28),
      size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 1 — KAPAK
#  Mesaj: NetGuard nedir, kim için, ne kadar
# ══════════════════════════════════════════════════════════════════════════════
def s01_kapak(prs):
    s = blank(prs); bg(s)

    # Sol lacivert alan
    box(s, 0, 0, Inches(5.0), SLIDE_H, fill=C_NAVY)
    box(s, Inches(5.0), 0, Inches(0.06), SLIDE_H, fill=C_BLUE)

    # Logo
    t(s, "Net",   Inches(0.5), Inches(1.6), Inches(4.3), Inches(1.1),
      size=80, bold=True, color=C_WHITE)
    t(s, "Guard", Inches(0.5), Inches(2.55), Inches(4.3), Inches(1.1),
      size=80, bold=True, color=C_BLUE)

    # Kısa açıklama
    line(s, Inches(0.5), Inches(3.7), Inches(4.1),
         color=RGBColor(0x33,0x41,0x55), thick=1.5)
    t(s, "Açık Kaynak\nAğ Güvenliği İzleme Platformu",
      Inches(0.5), Inches(3.85), Inches(4.2), Inches(0.9),
      size=16, color=RGBColor(0x94,0xA3,0xB8))

    # Kişi bilgisi
    t(s, "Mehmet Çapar", Inches(0.5), Inches(5.6), Inches(4.2), Inches(0.38),
      size=14, bold=True, color=C_WHITE)
    t(s, "Sakarya Üniversitesi — Bilgisayar Mühendisliği — 2026",
      Inches(0.5), Inches(5.95), Inches(4.2), Inches(0.28),
      size=11, color=RGBColor(0x64,0x74,0x8B))
    t(s, "Danışman: Prof. Dr. İbrahim ÖZÇELİK",
      Inches(0.5), Inches(6.22), Inches(4.2), Inches(0.28),
      size=11, italic=True, color=RGBColor(0x64,0x74,0x8B))

    # Sağ — tek büyük tagline
    t(s, "\"Splunk $50K/yıl.\nQRadar $30K/yıl.\nNetGuard: Ücretsiz,\n30 dakikada kurulum.\"",
      Inches(5.3), Inches(0.9), Inches(7.7), Inches(2.2),
      size=24, bold=True, color=C_NAVY, italic=True)

    # 4 KPI
    kpis = [
        ("83K",    "satır kod",  C_BLUE,   C_BLUE_LT),
        ("2.788",  "test / 0 hata", C_GREEN, C_GREEN_LT),
        ("50+",    "sigma kuralı", C_VIOLET, C_VIO_LT),
        ("5.9 sn", "ort. tespit", C_AMBER,  C_AMBER_LT),
    ]
    for i, (v, l, ac, bgc) in enumerate(kpis):
        kpi(s, Inches(5.3) + i * Inches(2.0), Inches(3.4),
            Inches(1.85), Inches(1.0), v, l, ac, bgc)

    # Standartlar
    box(s, Inches(5.3), Inches(4.6), Inches(7.7), Inches(0.42), fill=C_PANEL)
    t(s, "NIST SP 800-94  ·  CIS Controls v8.1  ·  MITRE ATT&CK v17  ·  SANS NSM",
      Inches(5.3), Inches(4.65), Inches(7.7), Inches(0.3),
      size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    photo_box(s, Inches(5.3), Inches(5.2), Inches(7.7), Inches(2.1),
              "NetGuard Dashboard — Overview Ekranı")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 2 — PROBLEM
#  Mesaj: Orta ölçekli şirketler hem tehdit altında hem bütçesiz
# ══════════════════════════════════════════════════════════════════════════════
def s02_problem(prs):
    s = blank(prs); bg(s)
    header(s, "Şirketler tehdit altında — ama bütçeleri yok", accent=C_RED)

    # 2 büyük stat — sol
    for i, (val, lbl, note, col, bgc) in enumerate([
        ("197 GÜN",
         "Ortalama saldırı tespit süresi",
         "Bu sürede saldırgan ağda serbestçe dolaşıyor",
         C_RED, C_RED_LT),
        ("$4.45 Milyon",
         "Ortalama veri ihlali maliyeti",
         "Tek bir ihlal küçük şirket için yıkıcı olabilir",
         C_AMBER, C_AMBER_LT),
    ]):
        cy = Inches(1.3) + i * Inches(2.5)
        box(s, Inches(0.5), cy, Inches(5.5), Inches(2.2),
            fill=bgc, line=col, lw=2.5)
        t(s, val, Inches(0.5), cy + Inches(0.15), Inches(5.5), Inches(1.0),
          size=52, bold=True, color=col, align=PP_ALIGN.CENTER)
        t(s, lbl, Inches(0.5), cy + Inches(1.1), Inches(5.5), Inches(0.4),
          size=15, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        line(s, Inches(0.9), cy + Inches(1.55), Inches(4.6), color=col, thick=0.75)
        t(s, note, Inches(0.5), cy + Inches(1.65), Inches(5.5), Inches(0.35),
          size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Sağ — maliyet karşılaştırması
    box(s, Inches(6.3), Inches(1.3), Inches(6.6), Inches(5.6),
        fill=C_PANEL, line=C_BORDER, lw=0.75)
    t(s, "Mevcut çözümler çok pahalı",
      Inches(6.5), Inches(1.4), Inches(6.3), Inches(0.4),
      size=15, bold=True, color=C_NAVY)

    costs = [
        ("Darktrace NDR",    100, C_RED,    "~$100K/yıl"),
        ("Splunk SIEM",       50, C_AMBER,  "~$50K/yıl"),
        ("IBM QRadar",        30, C_AMBER,  "~$30K/yıl"),
        ("Security Onion",     3, C_GREEN,  "Ücretsiz — günler süren kurulum"),
        ("NetGuard  ✓",        3, C_BLUE,   "Ücretsiz — 30 dakika"),
    ]
    bw = Inches(4.2)
    for i, (name, pct, col, label) in enumerate(costs):
        by = Inches(1.95) + i * Inches(0.82)
        bold_row = name.startswith("NetGuard")
        t(s, name, Inches(6.5), by, Inches(2.1), Inches(0.3),
          size=12, bold=bold_row, color=col if bold_row else C_DARK)
        bar_w = Inches(0.15) + bw * pct / 100
        box(s, Inches(6.5), by + Inches(0.34), bar_w, Inches(0.24), fill=col)
        t(s, label, Inches(6.5) + bar_w + Inches(0.1), by + Inches(0.32),
          Inches(2.5), Inches(0.3), size=11, color=col, bold=bold_row)

    bottom_bar(s, "Kaynak: Verizon DBIR 2025  ·  IBM Cost of Data Breach 2025")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 3 — ÇÖZÜM
#  Mesaj: NetGuard 3 şey yapar — Topla, Tespit Et, Engelle
# ══════════════════════════════════════════════════════════════════════════════
def s03_cozum(prs):
    s = blank(prs); bg(s)
    header(s, "NetGuard: Topla → Tespit Et → Engelle", accent=C_BLUE)

    # 3 büyük adım kutusu
    steps = [
        ("1\nTOPLA",       C_BLUE,   C_BLUE_LT,
         "9 farklı kaynaktan\nağ trafiği toplanır",
         ["Zeek, Suricata, NetFlow", "Windows Agent, Syslog", "OpenCanary, SNMP"]),
        ("2\nTESPİT ET",   C_VIOLET, C_VIO_LT,
         "4 farklı motor\nes zamanlı analiz eder",
         ["İmza tabanlı (Sigma)", "Kill chain korelasyonu", "Anomali + Beaconing"]),
        ("3\nUYAR & RAPORLA", C_GREEN, C_GREEN_LT,
         "Tehdit tespit edilince\nanında bildirim gönderilir",
         ["E-posta + webhook uyarısı", "Dashboard görselleştirme", "Audit log kaydı"]),
    ]

    for i, (title, col, bgc, subtitle, items) in enumerate(steps):
        cx = Inches(0.5) + i * Inches(4.3)
        # Ana kutu
        box(s, cx, Inches(1.3), Inches(3.9), Inches(5.5),
            fill=bgc, line=col, lw=2.5)
        # Numara dairesi
        box(s, cx + Inches(1.55), Inches(1.5), Inches(0.8), Inches(0.8), fill=col)
        t(s, str(i+1), cx + Inches(1.55), Inches(1.5), Inches(0.8), Inches(0.8),
          size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Başlık
        t(s, title.split("\n")[1], cx, Inches(2.45), Inches(3.9), Inches(0.55),
          size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        # Alt başlık
        t(s, subtitle, cx, Inches(3.0), Inches(3.9), Inches(0.7),
          size=13, color=C_DARK, align=PP_ALIGN.CENTER)
        line(s, cx + Inches(0.4), Inches(3.75), Inches(3.1), color=col, thick=1.0)
        # 3 madde
        for j, item in enumerate(items):
            jy = Inches(3.9) + j * Inches(0.5)
            box(s, cx + Inches(0.4), jy + Inches(0.1), Pt(8), Pt(8), fill=col)
            t(s, item, cx + Inches(0.58), jy, Inches(3.2), Inches(0.38),
              size=13, color=C_DARK)
        # Ok (son hariç)
        if i < 2:
            t(s, "→", cx + Inches(3.92), Inches(3.9), Inches(0.42), Inches(0.5),
              size=28, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    bottom_bar(s, "50–500 çalışanlı şirketler için  ·  Sıfır lisans maliyeti  ·  Docker ile 30 dakika kurulum")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 4 — MİMARİ
#  Mesaj: Kaynak → tek tablo → motor → yanıt — basit pipeline
# ══════════════════════════════════════════════════════════════════════════════
def s04_mimari(prs):
    s = blank(prs); bg(s)
    header(s, "Sistem Mimarisi — 4 Katmanlı Pipeline", accent=C_BLUE)

    phases = [
        ("VERİ\nKAYNAKLARI", C_BLUE,
         ["9 farklı kaynak", "Zeek · Suricata · NetFlow",
          "Windows · Syslog · SNMP", "OpenCanary Honeypot"]),
        ("ORTAK\nTABLO", C_VIOLET,
         ["normalized_logs", "ECS uyumlu format",
          "PostgreSQL + TimescaleDB", "Community ID pivot"]),
        ("TESPİT\nMOTORU", C_AMBER,
         ["pySigma 50+ kural", "Kill chain (5 aşama)",
          "Anomali (IsolationForest)", "Beaconing (IAT)"]),
        ("UYARI &\nRAPORLAMA", C_GREEN,
         ["E-posta + webhook", "Audit log",
          "AI açıklama", "Dashboard"]),
    ]

    pw = Inches(2.85)
    ph = Inches(5.3)

    for i, (title, col, items) in enumerate(phases):
        cx = Inches(0.45) + i * Inches(3.2)
        box(s, cx, Inches(1.25), pw, ph, fill=C_PANEL, line=col, lw=2.0)
        # Başlık bandı
        box(s, cx, Inches(1.25), pw, Inches(0.85), fill=col)
        t(s, title, cx, Inches(1.25), pw, Inches(0.85),
          size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Maddeler
        for j, item in enumerate(items):
            jy = Inches(2.25) + j * Inches(0.75)
            box(s, cx + Inches(0.25), jy + Inches(0.12), Pt(9), Pt(9), fill=col)
            t(s, item, cx + Inches(0.45), jy, pw - Inches(0.55), Inches(0.45),
              size=13, color=C_DARK)
        # Ok
        if i < 3:
            t(s, "→", cx + pw + Inches(0.15), Inches(3.55), Inches(0.35), Inches(0.5),
              size=26, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    bottom_bar(s, "FastAPI + Next.js 16 + Docker Compose 6 servis  ·  60 saniyelik korelasyon döngüsü")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 5 — VERİ KAYNAKLARI
#  Mesaj: 7 kaynaktan veri toplanıyor — 4×2 grid
# ══════════════════════════════════════════════════════════════════════════════
def s05_kaynaklar(prs):
    s = blank(prs); bg(s)
    header(s, "7 Veri Kaynağı — Tüm Ağ Görünür", accent=C_BLUE)

    sources = [
        ("Zeek NSM",      "14 log tipi\nconn · dns · http · ssl · ssh · smtp",  C_BLUE),
        ("Suricata IDS",  "45.000+ ET imzası\nHTTP · TLS · SSH anomali tespiti", C_RED),
        ("NetFlow",       "v5 / v9 / IPFIX / sFlow\nDoğu-batı trafik akışları", C_VIOLET),
        ("Windows Agent", "60+ olay ID, 5 kanal\nSecurity · Sysmon · PowerShell", C_AMBER),
        ("OpenCanary",    "12 sahte servis\nHoneypot tuzağı",                    C_GREEN),
        ("Syslog & SNMP", "UDP 514 · v2c/v3\nGüvenlik duvarı · cihaz metrikleri",C_MUTED),
        ("Community ID",  "RFC standardı korelasyon\nZeek ↔ Suricata ↔ NetFlow", C_BLUE),
    ]

    # 4+3 düzeni: ilk satır 4 kart, ikinci satır 3 kart ortalı
    cw = Inches(3.05)
    ch = Inches(2.2)
    gx = Inches(0.17)
    gy = Inches(0.25)

    layout = [
        # (source_index, row, col_offset)
        (0, 0, 0), (1, 0, 1), (2, 0, 2), (3, 0, 3),
        (4, 1, 0), (5, 1, 1), (6, 1, 2),
    ]
    # İkinci satırı ortala: 3 kart, ortada bırak
    row1_x_start = Inches(0.45)
    row2_total_w = 3 * cw + 2 * gx
    row2_x_start = (SLIDE_W - row2_total_w) / 2

    for src_i, row, col in layout:
        name, desc, col_c = sources[src_i]
        if row == 0:
            cx = row1_x_start + col * (cw + gx)
        else:
            cx = row2_x_start + col * (cw + gx)
        cy = Inches(1.25) + row * (ch + gy)

        box(s, cx, cy, cw, ch, fill=C_PANEL, line=col_c, lw=1.5)
        box(s, cx, cy, Inches(0.07), ch, fill=col_c)
        t(s, name, cx + Inches(0.17), cy + Inches(0.1),
          cw - Inches(0.22), Inches(0.38),
          size=13, bold=True, color=col_c)
        t(s, desc, cx + Inches(0.17), cy + Inches(0.55),
          cw - Inches(0.22), Inches(1.45),
          size=12, color=C_DARK)

    bottom_bar(s, "Tüm kaynaklar → normalized_logs (ECS uyumlu merkezi tablo) → PostgreSQL 16 + TimescaleDB")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 6 — TESPİT MOTORU
#  Mesaj: 4 motor eş zamanlı çalışır — hiçbir saldırı kaçmaz
# ══════════════════════════════════════════════════════════════════════════════
def s06_tespit(prs):
    s = blank(prs); bg(s)
    header(s, "4 Tespit Motoru — Eş Zamanlı Çalışır", accent=C_VIOLET)

    engines = [
        ("İmza\nTabanlı", C_BLUE, C_BLUE_LT,
         "50+ Sigma Kuralı",
         "Bilinen saldırı kalıpları\nanında tespit edilir",
         ["Port taraması", "Brute force", "DNS tünelleme"]),
        ("Kill Chain\nKorelasyon", C_VIOLET, C_VIO_LT,
         "5 Aşama / 30 dk",
         "Dağınık olaylar\nbir saldırı zinciri olarak görülür",
         ["RECON → WEAPONIZE", "ACCESS → LATERAL", "Tam zincir = blokaj"]),
        ("Davranış\nAnomalisi", C_AMBER, C_AMBER_LT,
         "Makine Öğrenmesi",
         "Normal olmayan trafik\notomatik işaretlenir",
         ["Trafik hacmi spike", "Yeni port / protokol", "7 günlük baseline"]),
        ("C2 Beaconing\nTespiti", C_GREEN, C_GREEN_LT,
         "IAT Analizi",
         "Zararlı yazılımın\nkomuta merkezi bağlantısı",
         ["Düzenli aralıklı trafik", "JA4 TLS parmak izi", "Tehdit puanı 0–100"]),
    ]

    for i, (title, col, bgc, badge, subtitle, items) in enumerate(engines):
        cx = Inches(0.45) + i * Inches(3.25)
        box(s, cx, Inches(1.25), Inches(3.0), Inches(5.65),
            fill=bgc, line=col, lw=2.0)
        # Numara
        box(s, cx + Inches(1.1), Inches(1.4), Inches(0.8), Inches(0.8), fill=col)
        t(s, str(i+1), cx + Inches(1.1), Inches(1.4), Inches(0.8), Inches(0.8),
          size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Başlık
        t(s, title, cx, Inches(2.3), Inches(3.0), Inches(0.6),
          size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        # Badge
        box(s, cx + Inches(0.55), Inches(2.95), Inches(1.9), Inches(0.28), fill=col)
        t(s, badge, cx + Inches(0.55), Inches(2.95), Inches(1.9), Inches(0.28),
          size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Açıklama
        t(s, subtitle, cx, Inches(3.35), Inches(3.0), Inches(0.6),
          size=12, color=C_DARK, align=PP_ALIGN.CENTER)
        line(s, cx + Inches(0.3), Inches(4.0), Inches(2.4), color=col, thick=0.75)
        # Maddeler
        for j, item in enumerate(items):
            jy = Inches(4.1) + j * Inches(0.5)
            box(s, cx + Inches(0.3), jy + Inches(0.1), Pt(8), Pt(8), fill=col)
            t(s, item, cx + Inches(0.47), jy, Inches(2.4), Inches(0.38),
              size=12, color=C_DARK)

    bottom_bar(s, "MITRE ATT&CK v17 eşleştirmesi  ·  Kill chain tespiti <60 saniye  ·  2.788 test — 0 hata")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 7 — CANLI DEMO
#  Mesaj: Gerçek saldırı 59 saniyede tespit edildi
# ══════════════════════════════════════════════════════════════════════════════
def s07_demo(prs):
    s = blank(prs); bg(s)
    header(s, "Gerçek Saldırı Testi — 59 Saniyede Tespit",
           "Kali Linux → NetGuard Server  |  10 Haziran 2026", accent=C_RED)

    # Büyük karşılaştırma
    box(s, Inches(0.45), Inches(1.25), Inches(4.0), Inches(1.5),
        fill=C_RED_LT, line=C_RED, lw=2.0)
    t(s, "197 GÜN", Inches(0.45), Inches(1.35), Inches(4.0), Inches(0.75),
      size=44, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    t(s, "Endüstri ortalaması tespit süresi",
      Inches(0.45), Inches(2.05), Inches(4.0), Inches(0.35),
      size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    t(s, "→", Inches(4.55), Inches(1.7), Inches(0.6), Inches(0.6),
      size=32, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    box(s, Inches(5.25), Inches(1.25), Inches(4.0), Inches(1.5),
        fill=C_GREEN_LT, line=C_GREEN, lw=2.5)
    t(s, "59 SANİYE", Inches(5.25), Inches(1.35), Inches(4.0), Inches(0.75),
      size=44, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    t(s, "NetGuard kill chain tespit süresi",
      Inches(5.25), Inches(2.05), Inches(4.0), Inches(0.35),
      size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    t(s, "34× DAHA HIZLI",
      Inches(9.4), Inches(1.5), Inches(3.5), Inches(0.6),
      size=22, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # 5 aşamalı saldırı şeridi
    stages = [
        ("RECON",       "Port taraması\n(nmap)",        C_AMBER),
        ("WEAPONIZE",   "Brute force\n(660 deneme)",    C_VIOLET),
        ("ACCESS",      "SSH girişi\nbaşarılı",         C_RED),
        ("LATERAL",     "İç ağ\ntaraması",              C_RED),
        ("FULL\nATTACK", "Zincir tamam\nEmail uyarısı!", C_RED),
    ]

    sw = Inches(12.4) / 5
    # Çizgi
    box(s, Inches(0.45), Inches(3.52), Inches(12.4), Pt(2),
        fill=C_PANEL2)

    for i, (stage, note, col) in enumerate(stages):
        cx = Inches(0.45) + i * sw
        # Aşama kutusu
        box(s, cx + Inches(0.1), Inches(2.9), sw - Inches(0.2), Inches(0.55), fill=col)
        t(s, stage, cx + Inches(0.1), Inches(2.9), sw - Inches(0.2), Inches(0.55),
          size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Nokta
        box(s, cx + sw/2 - Inches(0.12), Inches(3.39),
            Inches(0.24), Inches(0.24), fill=col)
        # Not
        box(s, cx + Inches(0.1), Inches(3.8), sw - Inches(0.2), Inches(0.7),
            fill=C_PANEL, line=col, lw=0.75)
        t(s, note, cx + Inches(0.1), Inches(3.83), sw - Inches(0.2), Inches(0.65),
          size=11, color=C_DARK, align=PP_ALIGN.CENTER)
        if i < 4:
            t(s, "→", cx + sw - Inches(0.12), Inches(2.98), Inches(0.25), Inches(0.38),
              size=16, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Ekran görüntüsü
    photo_box(s, Inches(0.45), Inches(4.7), Inches(12.4), Inches(2.4),
              "Kill Chain Timeline — Slaytlara eklenecek ekran görüntüsü")

    bottom_bar(s, "GNS3 Lab: Kali Linux · Windows Server 2022 (WIN-9DUCSU7LDJ0) · NetGuard VM 192.168.203.134")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 8 — DASHBOARD
#  Mesaj: 37 sayfalık modern arayüz — her şey tek ekranda
# ══════════════════════════════════════════════════════════════════════════════
def s08_dashboard(prs):
    s = blank(prs); bg(s)
    header(s, "37 Sayfalık Modern Dashboard",
           "Gerçek zamanlı görünürlük — tek ekranda", accent=C_BLUE)

    # Büyük ekran görüntüsü
    photo_box(s, Inches(0.45), Inches(1.25), Inches(8.0), Inches(5.1),
              "NetGuard Dashboard — Overview")

    # Sağ — 6 öne çıkan özellik
    features = [
        ("Kill Chain Timeline",  "Saldırı zinciri görsel olarak izlenir",          C_RED),
        ("AI Alert Explainer",   "Her uyarıyı Llama 3.3 70B açıklar",             C_VIOLET),
        ("Threat Hunt",          "Kod yazmadan tehdit araştırması",                 C_BLUE),
        ("MITRE ATT&CK",         "Kapsam analizi ve boşluk haritası",              C_AMBER),
        ("Live Log Stream",      "WebSocket ile anlık log akışı",                  C_GREEN),
        ("Sigma Rule Wizard",    "4 adımda yeni tespit kuralı oluşturma",          C_VIOLET),
    ]

    for i, (title, desc, col) in enumerate(features):
        fy = Inches(1.25) + i * Inches(0.85)
        box(s, Inches(8.65), fy, Inches(4.3), Inches(0.75),
            fill=C_PANEL, line=col, lw=1.5)
        box(s, Inches(8.65), fy, Inches(0.07), Inches(0.75), fill=col)
        t(s, title, Inches(8.82), fy + Inches(0.06),
          Inches(4.1), Inches(0.3), size=13, bold=True, color=col)
        t(s, desc, Inches(8.82), fy + Inches(0.38),
          Inches(4.1), Inches(0.28), size=11, color=C_MUTED)

    bottom_bar(s, "React 19 · Next.js 16 · ECharts · WebSocket · TanStack Query · Zustand")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 9 — KULLANILAN TEKNOLOJİLER
#  Mesaj: Modern, kanıtlanmış teknoloji yığını — 6 kategori
# ══════════════════════════════════════════════════════════════════════════════
def s09_teknolojiler(prs):
    s = blank(prs); bg(s)
    header(s, "Kullanılan Teknolojiler",
           "Modern, açık kaynak ve endüstri standardı araçlar", accent=C_BLUE)

    # 6 kategori — 3 sütun × 2 satır
    categories = [
        {
            "title": "Frontend",
            "color": C_BLUE,
            "bg":    C_BLUE_LT,
            "techs": [
                ("React 19",        "UI bileşen kütüphanesi"),
                ("Next.js 16",      "SSR + App Router"),
                ("TypeScript",      "Tip güvenli geliştirme"),
                ("Tailwind CSS",    "Utility-first stil"),
                ("ECharts",         "Grafik & görselleştirme"),
                ("TanStack Query",  "Sunucu durumu yönetimi"),
            ],
        },
        {
            "title": "Backend",
            "color": C_VIOLET,
            "bg":    C_VIO_LT,
            "techs": [
                ("Python 3.12",     "Ana programlama dili"),
                ("FastAPI",         "Async REST API çerçevesi"),
                ("pySigma v2",      "Sigma kural motoru"),
                ("scikit-learn",    "IsolationForest anomali"),
                ("pyotp",           "TOTP / MFA"),
                ("SlowAPI",         "Rate limiting"),
            ],
        },
        {
            "title": "Veritabanı",
            "color": C_GREEN,
            "bg":    C_GREEN_LT,
            "techs": [
                ("PostgreSQL 16",   "İlişkisel veritabanı"),
                ("TimescaleDB",     "Zaman serisi uzantısı"),
                ("Alembic",         "Şema migrasyon aracı"),
                ("Row-Level Sec.",  "Multi-tenant izolasyon"),
                ("Hypertable",      "1 günlük bölümleme"),
                ("Sıkıştırma",      "7 günde ~8:1 oran"),
            ],
        },
        {
            "title": "NSM Katmanı",
            "color": C_RED,
            "bg":    C_RED_LT,
            "techs": [
                ("Zeek",            "14 log tipi — ağ analizi"),
                ("Suricata",        "45.000+ ET imza kuralı"),
                ("OpenCanary",      "Honeypot — 12 servis"),
                ("communityid",     "Zeek↔Suricata pivot"),
                ("pyshark",         "BPF / paket yakalama"),
                ("Sigma YAML",      "50+ tespit kuralı"),
            ],
        },
        {
            "title": "Altyapı",
            "color": C_AMBER,
            "bg":    C_AMBER_LT,
            "techs": [
                ("Docker Compose",  "6 servis tek komutla"),
                ("nginx",           "Ters proxy + statik"),
                ("systemd",         "Servis yönetimi"),
                ("Alembic 21",      "Tam migration geçmişi"),
                ("pytest",          "2.788 otomatik test"),
                ("GitHub",          "Sürüm kontrolü"),
            ],
        },
        {
            "title": "Entegrasyonlar",
            "color": C_MUTED,
            "bg":    C_PANEL,
            "techs": [
                ("Groq Llama 3.3",  "AI alert açıklama"),
                ("AbuseIPDB",       "Tehdit istihbaratı"),
                ("Feodo Tracker",   "C2 sunucu listesi"),
                ("ThreatFox",       "IOC veritabanı"),
                ("GreyNoise",       "Tarama / gürültü filtresi"),
                ("CISA KEV",        "Bilinen istismar listesi"),
            ],
        },
    ]

    cw = Inches(4.1)
    ch = Inches(2.75)
    gx = Inches(0.16)
    gy = Inches(0.18)

    for i, cat in enumerate(categories):
        row = i // 3
        ci  = i %  3
        cx = Inches(0.45) + ci * (cw + gx)
        cy = Inches(1.2)  + row * (ch + gy)
        col = cat["color"]

        # Kart arka planı
        box(s, cx, cy, cw, ch, fill=cat["bg"], line=col, lw=1.5)
        # Başlık bandı
        box(s, cx, cy, cw, Inches(0.4), fill=col)
        t(s, cat["title"], cx, cy, cw, Inches(0.4),
          size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Teknoloji satırları — 2 sütun × 3 satır
        tw = (cw - Inches(0.2)) / 2
        for j, (name, desc) in enumerate(cat["techs"]):
            tc = j % 2        # sütun
            tr = j // 2       # satır
            tx = cx + Inches(0.1) + tc * tw
            ty = cy + Inches(0.48) + tr * Inches(0.72)
            box(s, tx, ty, tw - Inches(0.05), Inches(0.62),
                fill=C_WHITE, line=col, lw=0.5)
            t(s, name, tx + Inches(0.08), ty + Inches(0.04),
              tw - Inches(0.18), Inches(0.26),
              size=12, bold=True, color=col)
            t(s, desc, tx + Inches(0.08), ty + Inches(0.32),
              tw - Inches(0.18), Inches(0.24),
              size=10, color=C_MUTED)

    bottom_bar(s, "Tüm bileşenler açık kaynak  ·  83.000 satır kod  ·  Endüstri standardı araçlar")


# ══════════════════════════════════════════════════════════════════════════════
#  SLAYT 10 — SONUÇLAR & İLETİŞİM
#  Mesaj: NetGuard çalışıyor — rakamlar kanıtlıyor
# ══════════════════════════════════════════════════════════════════════════════
def s10_kapanis(prs):
    s = blank(prs); bg(s)
    header(s, "Sonuçlar", accent=C_NAVY)

    # 4 büyük KPI
    kpis = [
        ("83,000",  "Satır Kod",       C_BLUE,   C_BLUE_LT),
        ("2.788",   "Test — 0 Hata",   C_GREEN,  C_GREEN_LT),
        ("5.9 sn",  "Ort. Tespit",     C_AMBER,  C_AMBER_LT),
        ("~92/100", "NSM Kapsamı",     C_VIOLET, C_VIO_LT),
    ]
    for i, (v, l, ac, bgc) in enumerate(kpis):
        kpi(s, Inches(0.45) + i * Inches(3.25), Inches(1.25),
            Inches(3.0), Inches(1.3), v, l, ac, bgc)

    # Sol — özet mesajları
    box(s, Inches(0.45), Inches(2.8), Inches(7.9), Inches(3.75),
        fill=C_PANEL, line=C_BORDER, lw=0.75)
    t(s, "NetGuard ne sağlıyor?",
      Inches(0.6), Inches(2.9), Inches(7.7), Inches(0.38),
      size=16, bold=True, color=C_NAVY)

    results = [
        (C_BLUE,   "Security Onion düzeyinde NSM — Docker ile 30 dakikada"),
        (C_GREEN,  "Sıfır lisans maliyeti — küçük bütçeli şirketler için erişilebilir"),
        (C_VIOLET, "4 tespit motoru eş zamanlı — saldırı zinciri 59 saniyede yakalandı"),
        (C_AMBER,  "Kurumsal güvenlik kontrolü — MFA, RLS, SHA-256 audit log"),
        (C_RED,    "2.788 otomatik test — 0 hata, üretim ortamında kanıtlandı"),
    ]
    for i, (col, text) in enumerate(results):
        ry = Inches(3.4) + i * Inches(0.6)
        box(s, Inches(0.6), ry + Inches(0.07), Pt(10), Pt(10), fill=col)
        t(s, text, Inches(0.8), ry, Inches(7.4), Inches(0.38),
          size=13, color=C_DARK)

    # Sağ — iletişim
    box(s, Inches(8.55), Inches(2.8), Inches(4.4), Inches(3.75),
        fill=C_NAVY, line=C_BLUE, lw=2.5)
    t(s, "Mehmet Çapar",
      Inches(8.55), Inches(3.2), Inches(4.4), Inches(0.55),
      size=24, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    line(s, Inches(8.85), Inches(3.83), Inches(3.8), color=C_BLUE, thick=1.5)
    t(s, "20mehmetcapar02@gmail.com",
      Inches(8.55), Inches(3.95), Inches(4.4), Inches(0.35),
      size=13, color=C_BLUE, align=PP_ALIGN.CENTER)
    t(s, "Sakarya Üniversitesi\nBilgisayar Mühendisliği",
      Inches(8.55), Inches(4.4), Inches(4.4), Inches(0.65),
      size=13, color=RGBColor(0x94,0xA3,0xB8), align=PP_ALIGN.CENTER)
    t(s, "Danışman:\nProf. Dr. İbrahim ÖZÇELİK",
      Inches(8.55), Inches(5.15), Inches(4.4), Inches(0.65),
      size=12, italic=True, color=RGBColor(0x64,0x74,0x8B),
      align=PP_ALIGN.CENTER)

    # Alt bant
    box(s, 0, SLIDE_H - Inches(0.5), SLIDE_W, Inches(0.5), fill=C_NAVY)
    t(s, "docker-compose up -d",
      Inches(0.5), SLIDE_H - Inches(0.45), Inches(3.5), Inches(0.38),
      size=16, bold=True, color=RGBColor(0x4A,0xDE,0x80))
    t(s, "·  2.788 test, 0 hata  ·  NSM kapsamı ~92/100  ·  30 dakika kurulum",
      Inches(3.9), SLIDE_H - Inches(0.45), Inches(9.0), Inches(0.38),
      size=14, color=C_WHITE)


# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    agenda = [
        (s01_kapak,    "Kapak"),
        (s02_problem,  "Problem"),
        (s03_cozum,    "Çözüm: Topla → Tespit → Engelle"),
        (s04_mimari,   "Mimari"),
        (s05_kaynaklar,"Veri Kaynakları"),
        (s06_tespit,   "Tespit Motoru"),
        (s07_demo,     "Canlı Demo"),
        (s08_dashboard,"Dashboard"),
        (s09_teknolojiler, "Kullanılan Teknolojiler"),
        (s10_kapanis,  "Sonuçlar & İletişim"),
    ]
    for i, (fn, name) in enumerate(agenda):
        print(f"Slayt {i+1}/10 — {name}...")
        fn(prs)
    prs.save(DEST)
    print(f"\nKaydedildi: {DEST}")

if __name__ == "__main__":
    main()
