"""
NetGuard — Poster üretici
Poster_Formati2.pptx şablonunu NetGuard içeriğiyle doldurur.
Çalıştır: python3 uret_poster.py
"""

import copy
import shutil
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

SRC  = "/home/mehmet/İndirilenler/Poster_Formati2.pptx"
DEST = "/home/mehmet/netguard/rapor/NetGuard_Poster.pptx"

# ── İçerik ────────────────────────────────────────────────────────────────────

TITLE = (
    "NetGuard: Orta Ölçekli İşletmeler için Açık Kaynak "
    "Ağ Güvenliği İzleme Platformu"
)

AUTHOR_LEFT = [
    ("Mehmet Çapar", True, True),          # (metin, bold, italic)
    ("Bilgisayar Mühendisliği Bölümü", False, False),
    ("Sakarya Üniversitesi, BBIF", False, False),
    ("20mehmetcapar02@gmail.com", False, False),
]

AUTHOR_RIGHT = [
    ("Danışman: Prof. Dr. İbrahim ÖZÇELİK", False, True),
    ("Bilgisayar Mühendisliği Bölümü", False, False),
    ("Sakarya Üniversitesi", False, False),
]

# SOL SÜTUN
GIRIS_BASLIK = "Giriş & Motivasyon"
GIRIS_METIN = (
    "Verizon 2025 DBIR verilerine göre ortalama veri ihlali tespit "
    "süresi 197 gündür; bu süre içinde saldırganlar ağ içinde serbestçe "
    "hareket edebilmektedir. Ticari NSM/SIEM çözümlerinin lisans maliyetleri "
    "(Splunk ~50.000 $/yıl, QRadar ~30.000 $/yıl) orta ölçekli kurumlar için "
    "sürdürülemezdir. NetGuard, 50–500 çalışanlı şirketleri hedefleyen, "
    "Docker ile 30 dakikada kurulabilen, sıfır lisans maliyetli açık kaynak bir "
    "Ağ Güvenliği İzleme (NSM) platformudur. NIST SP 800-94, CIS Controls v8.1 "
    "ve MITRE ATT&CK v17 çerçevelerine uygun olarak tasarlanmıştır."
)

TABLO_CAPTION = "Tablo 1: Rakip NSM/SIEM sistemlerle karşılaştırmalı analiz."

VERI_BASLIK = "Veri Toplama Katmanı"
VERI_METIN = (
    "NetGuard 9 heterojen kaynaktan veri toplar ve ECS uyumlu "
    "normalized_logs tablosunda birleştirir: Syslog (güvenlik duvarı, IDS logları), "
    "SNMP v2c/v3 (cihaz metrikleri), Zeek (14 log tipi: conn, dns, http, ssl, ssh, "
    "smtp, ftp, rdp, kerberos, smb, dce_rpc, weird, dpd, files), Suricata EVE JSON "
    "(45.000+ ET kuralı), NetFlow v5/v9/IPFIX/sFlow (ağ akış verileri), "
    "Windows Agent (60+ EID; Security, Sysmon, PowerShell, System, AppLocker kanalları), "
    "OpenCanary Honeypot (12 sahte servis tipi), Microsoft 365 ve Google Workspace "
    "bulut toplayıcıları. Community ID (RFC) ile Zeek, Suricata ve NetFlow kayıtları "
    "aynı TCP bağlantısı üzerinden çapraz pivot edilebilir."
)

SEKIL1_CAP = "Şekil 1: NetGuard sistem mimarisi ve veri akış pipeline'ı."
SEKIL2_CAP = "Şekil 2: Dashboard genel bakış paneli (bento grid düzeni)."

MIMARI_BASLIK = "Sistem Mimarisi"
MIMARI_METIN = (
    "Üç katmanlı pipeline: Toplama → Normalleştirme → Tespit/Yanıt. "
    "FastAPI backend, Next.js 16 + React 19 dashboard, PostgreSQL 16 + TimescaleDB "
    "(1 günlük bölümleme, 7 günlük sıkıştırma, RPO ~24 sn). "
    "Korelasyon motoru 60 saniyelik döngüde çalışır; pySigma v2 kuralları "
    "doğrudan PostgreSQL sorgularına dönüştürülür. "
    "Docker Compose ile 6 servis tek komutla başlatılır."
)

SEKIL3_CAP = "Şekil 3: Kill chain swimlane timeline görünümü."
MIMARI_EK = (
    "Kill chain dedektörü kaynak IP bazında 5 aşamayı (RECON → WEAPONIZE → "
    "ACCESS → LATERAL → FULL_ATTACK_CHAIN) 30 dakikalık pencerede takip eder. "
    "FULL_ATTACK_CHAIN tespitinde OPNsense REST API veya VyOS SSH üzerinden "
    "otomatik IP blokajı tetiklenir."
)

GUVENLIK_BASLIK = "Güvenlik Özellikleri"
GUVENLIK_METIN = (
    "JWT + TOTP çok faktörlü kimlik doğrulama (pyotp, RFC 6238). "
    "SHA-256 zincirli tamper-proof denetim günlüğü (NIST SP 800-92 §3.2). "
    "PostgreSQL Row-Level Security ile multi-tenant veri izolasyonu. "
    "At-rest alan şifrelemesi (Fernet/AES-128-CBC). "
    "Sistematik rate limiting: 60 istek/dakika (SlowAPI, OWASP API Security 2023). "
    "Alembic 21 migration; 21 tablo, 11'i RLS korumalı."
)

SEKIL4_CAP = "Şekil 4: Aktif yanıt ve denetim günlüğü akış şeması."

# SAĞ SÜTUN
TESPIT_BASLIK = "Tespit Katmanı"
TESPIT_METIN = (
    "pySigma v2 (50+ kural, 14 dosya): Ağ taraması, kimlik doğrulama, "
    "C2, exfiltrasyon ve Windows/Zeek/Suricata kategorileri. "
    "JSON Korelasyon Motoru: 60 saniyelik pencerede eşik tabanlı gruplama. "
    "IsolationForest Anomali Tespiti: contamination=0.02, Welford z-skor + IQR. "
    "C2 Beaconing: IAT (Inter-Arrival Time) analizi, Bessel düzeltmeli standart "
    "sapma — RITA BlackHat 2018 metodolojisi. "
    "Tehdit İstihbaratı: AbuseIPDB + Feodo Tracker + ThreatFox + GreyNoise, "
    "composite skor 0–100. CISA KEV entegrasyonu ile bilinen istismar "
    "güvenlik açıkları anlık izlenir. TLS/JA4 parmak izi ile C2 çerçeve tespiti "
    "entropi eşiği 4.0 üzerinde çalışır."
)

TEST_BASLIK = "Test Ortamı & Bulgular"
TEST_METIN = (
    "GNS3 sanal laboratuvar ortamı: OPNsense 26.1.2, VyOS rolling, "
    "Alpine WebServer, Windows Server 2022 (WIN-9DUCSU7LDJ0, 5 kanal aktif). "
    "Doğrulanan saldırı senaryoları:"
)

SEKIL6_CAP = "Şekil 6: GNS3 laboratuvar ağ topolojisi."

TEST_EK = (
    "2788 birim, çapraz ve entegrasyon testi: 0 hata. "
    "Kill chain tespiti <60 sn gecikme. "
    "Otomatik IP blokajı başarı oranı: %100 (5/5 senaryo). "
    "Docker Compose kurulum süresi: ~30 dakika. "
    "TimescaleDB sıkıştırma oranı: ~8:1."
)

SONUC_BASLIK = "Sonuçlar"
SONUC_METIN = (
    "NetGuard, Security Onion ile karşılaştırılabilir NSM kapsamını Docker ile "
    "30 dakikada ve sıfır lisans maliyetiyle sunar. Çok katmanlı tespit "
    "(imza+korelasyon+anomali+kill chain+beaconing), 37 sayfalık modern React "
    "dashboard, gerçek zamanlı WebSocket log stream ve AI Alert Explainer "
    "(Groq Llama 3.3 70B) ile kurumsal düzeyde güvenlik görünürlüğü sağlanır. "
    "Gelecek çalışmalar: UEBA temeli (60+ gün verisi), Arkime tam PCAP "
    "entegrasyonu ve TheHive SOAR bağlantısı."
)

KAYNAKLAR_BASLIK = "Kaynaklar"
KAYNAKLAR_METIN = (
    "[1] NIST SP 800-94 Rev.1 — IDPS Kılavuzu\n"
    "[2] CIS Controls v8.1, Center for Internet Security\n"
    "[3] MITRE ATT&CK v17, MITRE Corporation\n"
    "[4] Verizon DBIR 2025, Verizon Enterprise\n"
    "[5] CrowdStrike 2025 Global Threat Report\n"
    "[6] SANS NSM Course, R. Bejtlich\n"
    "[7] Community ID Spec, Corelight (RFC)"
)

# Karşılaştırma tablosu (10 satır x 3 sütun)
TABLE_DATA = [
    ["Sistem", "Yıllık Maliyet", "Kurulum Süresi"],
    ["Splunk Enterprise", "~50.000 $", "Haftalar"],
    ["IBM QRadar", "~30.000 $", "Haftalar"],
    ["Darktrace NDR", "~100.000 $", "Günler"],
    ["Security Onion", "Ücretsiz", "Günler"],
    ["Wazuh (EDR)", "Ücretsiz", "Saatler"],
    ["Zeek (standalone)", "Ücretsiz", "Saatler"],
    ["Suricata (standalone)", "Ücretsiz", "Saatler"],
    ["NetGuard (bu çalışma)", "Ücretsiz", "~30 dakika"],
    ["✓ = NSM kapsamı", "✗ = Yok", "~ = Kısmi"],
]

# Test senaryoları tablosu (10 satır x 3 sütun — Shape[23])
# Bunu karşılaştırma tablosu yerine test senaryoları için kullanıyoruz
# Ama zaten karşılaştırma tablosu için TABLE_DATA yeterince iyi

# ── Yardımcı fonksiyonlar ─────────────────────────────────────────────────────

def set_tf_single(tf, text, bold=None, italic=None, size_pt=None):
    """Tek paragraf, tek run — mevcut formatı korur, metni değiştirir."""
    para = tf.paragraphs[0]
    # Fazla paragrafları sil
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    # Run'u ayarla
    if para.runs:
        run = para.runs[0]
    else:
        run = para.add_run()
    run.text = text
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    # Fazla run'ları sil
    while len(para.runs) > 1:
        r = para.runs[-1]._r
        r.getparent().remove(r)


def add_para(tf, text, bold=False, italic=False, size_pt=None, align=None):
    """tf'e yeni paragraf ekler — mevcut son paragrafın formatını miras alır."""
    from pptx.oxml.ns import qn
    import copy as _copy
    last_p = tf.paragraphs[-1]._p
    new_p = _copy.deepcopy(last_p)
    # run metni temizle
    for r in new_p.findall(qn('a:r')):
        new_p.remove(r)
    last_p.getparent().append(new_p)
    para = tf.paragraphs[-1]
    run = para.add_run()
    run.text = text
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if size_pt:
        run.font.size = Pt(size_pt)
    if align:
        para.alignment = align


def set_tf_section(tf, header, body, header_size=25, body_size=20):
    """Bölüm başlığı + gövde metni yazar."""
    set_tf_single(tf, header, bold=False, italic=False, size_pt=header_size)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_para(tf, body, bold=False, italic=False, size_pt=body_size,
             align=PP_ALIGN.JUSTIFY)


def set_author_block(tf, lines):
    """[(metin, bold, italic), ...] listesini yazar."""
    set_tf_single(tf, lines[0][0], bold=lines[0][1], italic=lines[0][2])
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for text, bold, italic in lines[1:]:
        add_para(tf, text, bold=bold, italic=italic, align=PP_ALIGN.CENTER)


def set_caption(shape, text):
    tf = shape.text_frame
    set_tf_single(tf, text, bold=False, italic=True, size_pt=20)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def fill_table(tbl, data):
    """TABLE şeklini data listesiyle doldurur (satır x sütun)."""
    for r_idx, row_data in enumerate(data):
        if r_idx >= len(tbl.rows):
            break
        row = tbl.rows[r_idx]
        for c_idx, cell_text in enumerate(row_data):
            if c_idx >= len(row.cells):
                break
            cell = row.cells[c_idx]
            tf = cell.text_frame
            if tf.paragraphs:
                set_tf_single(tf, cell_text,
                              bold=(r_idx == 0),   # header satırı bold
                              italic=False, size_pt=18)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER


# ── Ana üretim ────────────────────────────────────────────────────────────────

shutil.copy2(SRC, DEST)
prs = Presentation(DEST)
slide = prs.slides[0]
shapes = list(slide.shapes)

# [0] Başlık
set_tf_single(shapes[0].text_frame, TITLE, bold=True, italic=False, size_pt=70)

# [1] Sol yazar
set_author_block(shapes[1].text_frame, AUTHOR_LEFT)

# [2] Sağ yazar
set_author_block(shapes[2].text_frame, AUTHOR_RIGHT)

# [3] Sol sütun: Giriş & Motivasyon
set_tf_section(shapes[3].text_frame, GIRIS_BASLIK, GIRIS_METIN)

# [4] Picture 2 — bırak (kullanıcı eklecek)

# [5] Çizelge caption
set_caption(shapes[5], TABLO_CAPTION)

# [6] Sol sütun: Veri Toplama Katmanı
set_tf_section(shapes[6].text_frame, VERI_BASLIK, VERI_METIN)

# [7][8] Picture 3 & 4 — bırak

# [9] Sol sütun: devam metni (veri toplama alt detay)
tf9 = shapes[9].text_frame
set_tf_single(tf9,
    "Pipeline: Kaynak → normalized_logs → Correlator → Kill Chain → Aktif Yanıt",
    bold=False, italic=False, size_pt=20)
tf9.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

# [10][11] Şekil 1 & 2 caption
set_caption(shapes[10], SEKIL1_CAP)
set_caption(shapes[11], SEKIL2_CAP)

# [12] Sol sütun: Sistem Mimarisi
set_tf_section(shapes[12].text_frame, MIMARI_BASLIK, MIMARI_METIN)

# [14] Resim 1 — bırak

# [15] Şekil 3 caption
set_caption(shapes[15], SEKIL3_CAP)

# [16] Mimari ek metin
tf16 = shapes[16].text_frame
set_tf_single(tf16, MIMARI_EK, bold=False, italic=False, size_pt=20)
tf16.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

# [17] Sol sütun: Güvenlik Özellikleri
tf17 = shapes[17].text_frame
set_tf_single(tf17, GUVENLIK_BASLIK, bold=True, italic=False, size_pt=25)
tf17.paragraphs[0].alignment = PP_ALIGN.CENTER
add_para(tf17, GUVENLIK_METIN, bold=False, italic=False, size_pt=20,
         align=PP_ALIGN.JUSTIFY)

# [18] Resim 3 — bırak

# [19] Şekil 4 caption
set_caption(shapes[19], SEKIL4_CAP)

# [20] Sağ sütun: Tespit Katmanı
set_tf_section(shapes[20].text_frame, TESPIT_BASLIK, TESPIT_METIN)

# [21] Picture 2 (sağ) — bırak

# [22] Sağ sütun: Test Ortamı & Bulgular
set_tf_section(shapes[22].text_frame, TEST_BASLIK, TEST_METIN)

# [23] TABLE — karşılaştırma tablosu
fill_table(shapes[23].table, TABLE_DATA)

# [24] Picture 4 (sağ) — bırak

# [25] Şekil 6 caption
set_caption(shapes[25], SEKIL6_CAP)

# [26] Test ek metin
tf26 = shapes[26].text_frame
set_tf_single(tf26, TEST_EK, bold=False, italic=False, size_pt=20)
tf26.paragraphs[0].alignment = PP_ALIGN.JUSTIFY

# [27] Sağ sütun: Sonuçlar
set_tf_section(shapes[27].text_frame, SONUC_BASLIK, SONUC_METIN)

# [28] Kaynaklar (AUTO_SHAPE)
tf28 = shapes[28].text_frame
set_tf_single(tf28, KAYNAKLAR_BASLIK, bold=True, italic=False, size_pt=25)
tf28.paragraphs[0].alignment = PP_ALIGN.CENTER
for line in KAYNAKLAR_METIN.split("\n"):
    add_para(tf28, line, bold=False, italic=False, size_pt=18,
             align=PP_ALIGN.LEFT)

# [29] Picture 5 (büyük) — bırak

prs.save(DEST)
print(f"Kaydedildi: {DEST}")
